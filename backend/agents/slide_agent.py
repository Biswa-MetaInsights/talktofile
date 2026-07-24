"""
Slide Agent — generates a PowerPoint presentation from document content.
Uses python-pptx to produce a .pptx file.
"""

import io
import json
from openai import AsyncOpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from core.session_store import DocumentData
from core.config import get_settings


_SYSTEM = """You are a presentation designer. Given document content, create a structured slide deck.

Generate 8-12 slides covering the key content. Each slide should have:
- A clear, concise title (max 8 words)
- 3-5 bullet points (each max 15 words)
- Optional speaker_note (1-2 sentences for the presenter)

The first slide is always a title slide with just a title and subtitle.
The last slide is always a summary/takeaways slide.

Return ONLY a JSON array (no markdown, no extra text):
[
  {"type": "title", "title": "...", "subtitle": "...", "speaker_note": "..."},
  {"type": "content", "title": "...", "bullets": ["...", "..."], "speaker_note": "..."},
  ...
]"""


_REFINE_SYSTEM = """You are a presentation editor. You are given an EXISTING slide deck (as a JSON
array) and an instruction from the user describing how to change it. Apply the instruction and
return the FULL updated deck.

Rules:
- Keep the same JSON structure: each slide is an object with keys "type" ("title" or "content"),
  "title", optional "subtitle" (title slides), "bullets" (array, content slides), and optional
  "speaker_note".
- Preserve slides and content the user did not ask to change. Only modify what the instruction
  implies (rewording, adding/removing/reordering slides or bullets, shortening, tone, etc.).
- Keep the first slide a "title" slide.
- Titles stay concise (max ~8 words); bullets stay concise (max ~15 words).

Return ONLY the updated JSON array — no markdown, no commentary."""


def _build_context(documents: list[DocumentData]) -> str:
    parts = []
    for doc in documents:
        if isinstance(doc.summary, dict):
            s = doc.summary
            parts.append(
                f"=== {doc.filename} ===\n"
                f"Overview: {s.get('overview','')}\n"
                f"Key points: {'; '.join(s.get('key_points', []))}\n"
                f"Topics: {', '.join(s.get('topics', []))}"
            )
        if doc.chunks:
            content = "\n\n".join(doc.chunks)[:8000]
            parts.append(f"Content:\n{content}")
    return "\n\n".join(parts)


def _parse_slides(raw: str) -> list[dict]:
    raw = (raw or "[]").strip()
    if raw.startswith("```"):
        raw = raw.split("\n", 1)[-1].rsplit("```", 1)[0].strip()
    try:
        slides = json.loads(raw)
        return slides if isinstance(slides, list) else []
    except json.JSONDecodeError:
        return []


async def generate_slides_data(documents: list[DocumentData]) -> list[dict]:
    settings = get_settings()
    client = AsyncOpenAI(api_key=settings.openai_api_key)
    context = _build_context(documents)

    response = await client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": _SYSTEM},
            {"role": "user", "content": f"Document content:\n\n{context}"},
        ],
        temperature=0.4,
        max_tokens=3000,
    )

    return _parse_slides(response.choices[0].message.content)


async def refine_slides_data(
    documents: list[DocumentData], current_slides: list[dict], instruction: str
) -> list[dict]:
    """Apply a natural-language instruction to an existing deck and return the full
    updated deck. The source document is provided as context so the model can add
    genuinely new, grounded slides/bullets when asked."""
    settings = get_settings()
    client = AsyncOpenAI(api_key=settings.openai_api_key)
    context = _build_context(documents)

    user_msg = (
        f"Source document (for grounding any additions):\n\n{context}\n\n"
        f"Existing slide deck (JSON):\n{json.dumps(current_slides, ensure_ascii=False)}\n\n"
        f"Instruction: {instruction}"
    )

    response = await client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": _REFINE_SYSTEM},
            {"role": "user", "content": user_msg},
        ],
        temperature=0.4,
        max_tokens=3000,
    )

    refined = _parse_slides(response.choices[0].message.content)
    # Never wipe the deck on a bad parse — fall back to what we had.
    return refined or current_slides


# ── Theming ───────────────────────────────────────────────────────────────────
# The frontend renders an HTML preview that mirrors these presets (see SlidesView
# SlideCanvas). Keep the two in sync when changing colours/layout.

_DEFAULT_ACCENT = RGBColor(0xE2, 0x61, 0x1B)  # brand orange
_DARK = RGBColor(0x30, 0x30, 0x30)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
_LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
_MUTED = RGBColor(0x6B, 0x72, 0x80)

_PRESETS = {"classic", "minimal", "bold"}


def _parse_accent(value) -> RGBColor:
    """Accept a '#RRGGBB' string; fall back to the brand orange."""
    if isinstance(value, str):
        v = value.strip().lstrip("#")
        if len(v) == 3:
            v = "".join(c * 2 for c in v)
        if len(v) == 6:
            try:
                return RGBColor(int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))
            except ValueError:
                pass
    return _DEFAULT_ACCENT


def _tint(color: RGBColor, factor: float) -> RGBColor:
    """Lighten a colour toward white by `factor` (0..1)."""
    r, g, b = color[0], color[1], color[2]
    return RGBColor(
        int(r + (255 - r) * factor),
        int(g + (255 - g) * factor),
        int(b + (255 - b) * factor),
    )


def _set_text(tf, text: str, size_pt: int, bold: bool = False, color: RGBColor = None):
    tf.text = text
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            if color:
                run.font.color.rgb = color


def _fill_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def build_pptx(
    slides_data: list[dict],
    doc_title: str = "Document",
    theme: dict | None = None,
    author: str | None = None,
) -> bytes:
    theme = theme or {}
    preset = theme.get("preset", "classic")
    if preset not in _PRESETS:
        preset = "classic"
    accent = _parse_accent(theme.get("accent"))
    author = (author or "").strip()

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # completely blank

    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        slide_type = slide_data.get("type", "content")

        if slide_type == "title":
            _build_title_slide(slide, slide_data, doc_title, preset, accent, author)
        else:
            _build_content_slide(slide, slide_data, preset, accent)

        note = slide_data.get("speaker_note", "")
        if note:
            slide.notes_slide.notes_text_frame.text = note

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_title_slide(slide, slide_data, doc_title, preset, accent, author):
    title = slide_data.get("title", doc_title)
    subtitle = slide_data.get("subtitle", "")

    if preset == "minimal":
        _fill_bg(slide, _WHITE)
        _add_rect(slide, Inches(1), Inches(4.35), Inches(3), Inches(0.06), accent)
        title_color, sub_color, meta_color = _DARK, _MUTED, accent
    elif preset == "bold":
        _fill_bg(slide, _NEAR_BLACK)
        _add_rect(slide, Inches(0), Inches(0), Inches(0.35), Inches(7.5), accent)
        title_color, sub_color, meta_color = accent, _WHITE, _tint(accent, 0.3)
    else:  # classic
        _fill_bg(slide, accent)
        title_color, sub_color, meta_color = _WHITE, _tint(accent, 0.75), _tint(accent, 0.75)

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.33), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    _set_text(tf, title, 44, bold=True, color=title_color)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.6), Inches(11.33), Inches(1))
        tf2 = sub_box.text_frame
        tf2.word_wrap = True
        _set_text(tf2, subtitle, 24, bold=False, color=sub_color)

    if author:
        meta_box = slide.shapes.add_textbox(Inches(1), Inches(5.7), Inches(11.33), Inches(0.6))
        tf3 = meta_box.text_frame
        tf3.word_wrap = True
        _set_text(tf3, f"Created by {author}", 16, bold=False, color=meta_color)


def _build_content_slide(slide, slide_data, preset, accent):
    title = slide_data.get("title", "")
    bullets = slide_data.get("bullets", [])

    if preset == "minimal":
        _fill_bg(slide, _WHITE)
        _add_rect(slide, Inches(0.5), Inches(0.4), Inches(0.09), Inches(0.7), accent)
        title_left = Inches(0.75)
        bullet_color = _DARK
        marker = "—"
    elif preset == "bold":
        _fill_bg(slide, _LIGHT_GRAY)
        _add_rect(slide, Inches(0.5), Inches(1.15), Inches(2.2), Inches(0.06), accent)
        title_left = Inches(0.5)
        bullet_color = _DARK
        marker = "•"
    else:  # classic
        _fill_bg(slide, _WHITE)
        _add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), accent)
        title_left = Inches(0.5)
        bullet_color = _DARK
        marker = "•"

    title_box = slide.shapes.add_textbox(title_left, Inches(0.25), Inches(12.33), Inches(0.9))
    tf = title_box.text_frame
    tf.word_wrap = True
    _set_text(tf, title, 28, bold=True, color=_DARK)

    if bullets:
        bullet_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.6), Inches(5.5))
        tf2 = bullet_box.text_frame
        tf2.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p.space_before = Pt(8)
            # Coloured marker run + dark text run.
            run_marker = p.add_run()
            run_marker.text = f"{marker}  "
            run_marker.font.size = Pt(20)
            run_marker.font.bold = True
            run_marker.font.color.rgb = accent
            run_text = p.add_run()
            run_text.text = bullet
            run_text.font.size = Pt(20)
            run_text.font.color.rgb = bullet_color


async def generate_presentation(documents: list[DocumentData]) -> bytes:
    slides_data = await generate_slides_data(documents)
    doc_title = documents[0].filename.rsplit(".", 1)[0] if documents else "Document"
    return build_pptx(slides_data, doc_title)
